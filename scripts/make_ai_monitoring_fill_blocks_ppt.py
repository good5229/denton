from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from make_ai_monitoring_ppt import (
    BG,
    BLUE,
    GOLD,
    GREEN,
    INK,
    LINE,
    LIGHT,
    MUTED,
    PURPLE,
    RED,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    WHITE,
    add_circle,
    add_rect,
    add_text,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "ppt"


def add_center_card(slide, x, y, w, h, num, title, lines, accent) -> None:
    add_rect(slide, x, y, w, h, WHITE, LINE)
    add_circle(slide, x + 24, y + 22, 44, accent, num)
    add_text(slide, x + 84, y + 22, w - 110, 34, title, 16, accent, True)
    add_text(slide, x + 84, y + 70, w - 115, h - 92, "\n".join(lines), 12, MUTED, True)


def add_metric_card(slide, x, y, w, h, title, value, note, accent) -> None:
    add_rect(slide, x, y, w, h, WHITE, LINE)
    add_text(slide, x + 26, y + 18, w - 52, 28, title, 14, accent, True)
    add_text(slide, x + 26, y + 55, w - 52, 42, value, 22, INK, True)
    add_text(slide, x + 26, y + 105, w - 52, 30, note, 11, MUTED, True)


def add_small_chip(slide, x, y, w, label, accent) -> None:
    add_rect(slide, x, y, w, 34, LIGHT, LIGHT, False)
    add_text(slide, x + 10, y + 8, w - 20, 20, label, 10, accent, True, PP_ALIGN.CENTER)


def build_copy_ready_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1920, 1080, BG, BG, False, 0)
    add_text(slide, 80, 58, 900, 58, "AI 산업 모니터링 보조 블록", 28, INK, True)
    add_text(slide, 82, 124, 1280, 34, "현재 슬라이드의 중앙·하단 빈 공간에 복사해 넣을 수 있는 구현 중심 PPT 요소", 15, MUTED, True)

    add_text(slide, 90, 205, 620, 48, "중앙 A · 구현 아키텍처", 19, INK, True)
    center_cards = [
        ("1", "문서 수집·파싱", ["웹 스크래핑", "PDF 본문 추출", "기업·산업 메타데이터"], BLUE),
        ("2", "문장 정제", ["수치·중복 제거", "유효 문장 선별", "분기별 샘플 구성"], TEAL),
        ("3", "NLP 추론", ["BERT 감성분류", "Trigram 키워드", "이벤트 문장 필터"], PURPLE),
        ("4", "지표 팩토리", ["TBCI·TIEI·TEEI", "BC-factors", "산업 유사도"], RED),
    ]
    x = 90
    for card in center_cards:
        add_center_card(slide, x, 265, 415, 182, *card)
        x += 440

    add_text(slide, 90, 515, 620, 48, "중앙 B · 지표 생성 로직", 19, INK, True)
    add_rect(slide, 90, 570, 830, 195, WHITE, LINE)
    add_text(slide, 130, 604, 170, 30, "감성 지수", 15, PURPLE, True)
    add_text(slide, 310, 604, 520, 30, "TBCI = (긍정 문장 - 부정 문장) / 감성 문장", 14, INK, True)
    add_text(slide, 130, 662, 170, 30, "이벤트 영향", 15, BLUE, True)
    add_text(slide, 310, 662, 520, 30, "TIEI = 이벤트 키워드 포함 문장 비중", 14, INK, True)
    add_text(slide, 130, 720, 170, 30, "이벤트 평가", 15, RED, True)
    add_text(slide, 310, 720, 520, 30, "TEEI = 이벤트 문장 내 긍정·부정 논조", 14, INK, True)

    add_rect(slide, 980, 570, 850, 195, WHITE, LINE)
    add_text(slide, 1020, 604, 180, 30, "요인 추출", 15, GOLD, True)
    add_small_chip(slide, 1215, 600, 115, "관계언", GOLD)
    add_small_chip(slide, 1348, 600, 115, "요인", GOLD)
    add_small_chip(slide, 1481, 600, 115, "평가", GOLD)
    add_small_chip(slide, 1614, 600, 115, "Trigram", GOLD)
    add_text(slide, 1020, 668, 180, 30, "유사도", 15, GREEN, True)
    add_text(slide, 1215, 668, 500, 30, "공통요인 분포의 KL divergence 역수", 14, INK, True)
    add_text(slide, 1020, 724, 180, 30, "검증", 15, RED, True)
    add_text(slide, 1215, 724, 500, 30, "GDP·BSI·주가지수 시차상관 검증", 14, INK, True)

    add_text(slide, 90, 835, 620, 48, "하단 · 포트폴리오 성과 지표", 19, INK, True)
    metrics = [
        ("데이터", "12.8만", "증권사 기업 리포트", BLUE),
        ("정제", "145만", "유효 문장 샘플", TEAL),
        ("모델", "BERT", "긍정·부정·중립 분류", PURPLE),
        ("키워드", "Trigram", "산업 이슈 자동 추출", GOLD),
        ("검증", "0.91", "선행지수 상관 예시", GREEN),
    ]
    x = 90
    for metric in metrics:
        add_metric_card(slide, x, 890, 325, 140, *metric)
        x += 350


def build_variant_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1920, 1080, BG, BG, False, 0)
    add_text(slide, 80, 58, 900, 58, "AI 산업 모니터링 보조 블록 · 압축형", 28, INK, True)
    add_text(slide, 82, 124, 1280, 34, "원본 슬라이드 공간이 좁을 때 쓰기 좋은 낮은 높이의 중앙·하단 블록", 15, MUTED, True)

    add_rect(slide, 120, 230, 1680, 210, WHITE, LINE)
    add_text(slide, 160, 265, 220, 36, "Pipeline", 18, INK, True)
    compact = [
        ("수집", "리포트·기업·업종"),
        ("정제", "문장 샘플링"),
        ("분류", "BERT 감성"),
        ("추출", "Trigram 요인"),
        ("검증", "거시지표 비교"),
    ]
    x = 410
    for idx, (title, body) in enumerate(compact, start=1):
        add_circle(slide, x, 270, 42, [BLUE, TEAL, PURPLE, GOLD, RED][idx - 1], str(idx))
        add_text(slide, x + 58, 262, 140, 26, title, 13, [BLUE, TEAL, PURPLE, GOLD, RED][idx - 1], True)
        add_text(slide, x + 58, 298, 160, 28, body, 11, MUTED, True)
        x += 270

    add_rect(slide, 120, 525, 530, 225, WHITE, LINE)
    add_text(slide, 158, 560, 260, 34, "원문 근거 추적", 18, BLUE, True)
    add_text(slide, 158, 616, 420, 56, "지표 수치만 제시하지 않고, 해당 지표를 만든 원문 문장과 키워드를 함께 보관", 13, MUTED, True)
    add_text(slide, 158, 700, 420, 28, "분석가 설명 가능성 확보", 13, INK, True)

    add_rect(slide, 695, 525, 530, 225, WHITE, LINE)
    add_text(slide, 733, 560, 260, 34, "산업 간 네트워크", 18, GREEN, True)
    add_text(slide, 733, 616, 420, 56, "공통으로 언급되는 요인 분포를 비교해 산업 간 유사도와 전이 가능성을 시각화", 13, MUTED, True)
    add_text(slide, 733, 700, 420, 28, "연관 산업군 탐색", 13, INK, True)

    add_rect(slide, 1270, 525, 530, 225, WHITE, LINE)
    add_text(slide, 1308, 560, 260, 34, "운영 자동화", 18, PURPLE, True)
    add_text(slide, 1308, 616, 420, 56, "보고서 수집부터 지표 산출, 시각화까지 반복 실행 가능한 Python 파이프라인", 13, MUTED, True)
    add_text(slide, 1308, 700, 420, 28, "주기적 산업 모니터링", 13, INK, True)

    add_rect(slide, 120, 865, 1680, 95, WHITE, LINE)
    add_text(slide, 160, 892, 230, 30, "핵심 메시지", 16, RED, True)
    add_text(slide, 405, 892, 1260, 32, "숫자 전망치가 놓치는 애널리스트의 정성 판단을 문장 단위로 정량화해 산업별 경기 신호로 변환", 15, INK, True)


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_copy_ready_slide(prs)
    build_variant_slide(prs)
    out = OUT / "ai_industry_monitoring_fill_blocks_editable.pptx"
    prs.save(out)
    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
