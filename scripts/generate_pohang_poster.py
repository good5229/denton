from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pohang"
DATA = ROOT / "data" / "processed"
RAW = ROOT / "data" / "raw" / "phase42_pohang"
HIERARCHICAL_VALIDATION = DATA / "phase64_hierarchical_aggregate_validation" / "phase64_small_to_middle_aggregate_validation_detail.csv"
FINAL_ACCURACY_REGISTRY = DATA / "phase98_final_middle_industry_accuracy_registry" / "phase98_final_middle_industry_accuracy_registry.csv"
NO_WORSE_REFINEMENT = DATA / "phase105_no_worse_refinement_guardrail" / "phase105_no_worse_refinement_registry.csv"
PHASE114_REFINEMENT = DATA / "phase114_block_routed_refinement_audit" / "phase114_refined_registry.csv"
PHASE127_STRICT_REFINEMENT = DATA / "phase127_precision_comwel_after_phase114" / "phase127_strict_registry.csv"
PHASE128_FLASH = DATA / "phase128_vintage_flash_redesign" / "phase128_vintage_middle_flash_detail.csv"
PHASE217_REPORTING = DATA / "phase217_public_safe_candidate_rerank_audit" / "phase217_reranked_guarded_registry.csv"
W, H = 3508, 4967
M, GAP = 72, 20
BODY_W = W - 2 * M
COL_W = (BODY_W - 2 * GAP) / 3
SLIDE_W_IN, SLIDE_H_IN = 594 / 25.4, 841 / 25.4
SX, SY = SLIDE_W_IN / W, SLIDE_H_IN / H

NAVY, BLUE, SKY, PAGE = "073B67", "2B6F9F", "DCECF5", "EEF4F7"
INK, MUTED, GRID, WHITE = "14242E", "50636F", "C7D4DC", "FFFFFF"
TEAL, ORANGE, RED, GOLD, PALE = "147D78", "E06A3B", "B33B32", "B98724", "F5F8FA"
GREEN = "2D8B72"
FONT, FONT_BOLD = "NanumBarunGothic", "NanumBarunGothic Bold"
FONT_SCALE = .61
MIN_FONT_PT = 10.0


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def apply_phase128_flash(hv: pd.DataFrame, city: str) -> pd.DataFrame:
    """Replace rejected initial split with Q4+1M historical middle-industry flash."""
    if not PHASE128_FLASH.exists():
        hv["flash_eok"] = hv.initial_predicted_gva_eok
        hv["flash_error_eok"] = hv.initial_error_gva_eok
        hv["flash_error_rate_pct"] = hv.initial_error_rate_pct
        return hv
    flash = pd.read_csv(PHASE128_FLASH, dtype={"middle_code": str})
    flash["middle_code"] = flash.middle_code.str.zfill(2)
    flash = flash[
        flash.city.eq(city)
        & flash.vintage_id.eq("Q4_plus_1m")
        & flash.share_model_id.eq("historical_middle_split")
        & flash.allowed_track.eq("flash_candidate")
    ][["parent_code", "middle_code", "flash_predicted_gva_eok", "flash_error_gva_eok", "flash_error_rate_pct"]]
    hv["middle_code"] = hv.middle_code.astype(str).str.zfill(2)
    hv = hv.merge(flash, on=["parent_code", "middle_code"], how="left")
    hv["flash_eok"] = hv.flash_predicted_gva_eok.fillna(hv.initial_predicted_gva_eok)
    hv["flash_error_eok"] = hv.flash_error_gva_eok.fillna(hv.initial_error_gva_eok)
    hv["flash_error_rate_pct"] = hv.flash_error_rate_pct.fillna(hv.initial_error_rate_pct)
    return hv.drop(columns=["flash_predicted_gva_eok", "flash_error_gva_eok"], errors="ignore")


def load_reporting_hv(city: str) -> pd.DataFrame:
    """Load middle-industry GVA validation rows with the latest reporting columns.

    Phase217 is the current public-reporting contract.  Older precision columns
    remain diagnostic only because some of them can be worse than the flash
    estimate.
    """
    if PHASE217_REPORTING.exists():
        hv = pd.read_csv(PHASE217_REPORTING, dtype={"middle_code": str})
        hv = hv[hv.city.eq(city)].copy()
        hv["middle_code"] = hv.middle_code.astype(str).str.zfill(2)
        hv["middle_label"] = hv.middle_label.fillna(hv.middle_code.astype(str))
        hv["actual_eok"] = hv.actual_gva_eok
        hv["flash_eok"] = hv.flash_predicted_gva_eok
        hv["flash_error_eok"] = hv.flash_error_gva_eok
        hv["flash_error_rate_pct"] = hv.flash_error_rate_pct
        hv["refined_eok"] = hv.phase217_guarded_predicted_gva_eok
        hv["refined_error_eok"] = hv.phase217_guarded_error_gva_eok
        hv["refined_error_rate_pct"] = hv.phase217_guarded_error_rate_pct
        return hv

    hv_source = PHASE127_STRICT_REFINEMENT if PHASE127_STRICT_REFINEMENT.exists() else PHASE114_REFINEMENT if PHASE114_REFINEMENT.exists() else NO_WORSE_REFINEMENT if NO_WORSE_REFINEMENT.exists() else FINAL_ACCURACY_REGISTRY
    hv = pd.read_csv(hv_source, dtype={"middle_code": str})
    hv = hv[hv.city.eq(city)].copy()
    hv["middle_label"] = hv.middle_label.fillna(hv.middle_code.astype(str))
    hv["actual_eok"] = hv.actual_gva_eok
    hv["refined_eok"] = hv["phase127_strict_predicted_gva_eok"] if "phase127_strict_predicted_gva_eok" in hv.columns else hv["phase114_predicted_gva_eok"] if "phase114_predicted_gva_eok" in hv.columns else hv["no_worse_refined_predicted_gva_eok"] if "no_worse_refined_predicted_gva_eok" in hv.columns else hv.protected_predicted_gva_eok
    hv = apply_phase128_flash(hv, city)
    hv["refined_error_eok"] = hv["phase127_strict_error_gva_eok"] if "phase127_strict_error_gva_eok" in hv.columns else hv["phase114_error_gva_eok"] if "phase114_error_gva_eok" in hv.columns else hv["no_worse_refined_error_gva_eok"] if "no_worse_refined_error_gva_eok" in hv.columns else hv.protected_error_gva_eok
    hv["refined_error_rate_pct"] = hv["phase127_strict_error_rate_pct"] if "phase127_strict_error_rate_pct" in hv.columns else hv["phase114_error_rate_pct"] if "phase114_error_rate_pct" in hv.columns else hv["no_worse_refined_error_rate_pct"] if "no_worse_refined_error_rate_pct" in hv.columns else hv.protected_error_rate_pct
    return hv


def parent_letters(parent_section: str) -> list[str]:
    parent_section = str(parent_section)
    if parent_section == "ERS":
        return ["E", "R", "S"]
    if parent_section == "MN0":
        return ["M", "N"]
    return [parent_section[0]]


def table_label(value: str, max_chars: int = 12) -> str:
    value = str(value)
    if "\n" in value or len(value) <= max_chars:
        return value
    pivots = [" 및 ", "·", " ", ";"]
    for pivot in pivots:
        positions = [i + len(pivot) for i in range(len(value)) if value.startswith(pivot, i)]
        if not positions:
            continue
        cut = min(positions, key=lambda p: abs(p - len(value) / 2))
        if 4 <= cut <= len(value) - 4:
            return value[:cut].strip() + "\n" + value[cut:].strip()
    return value[:max_chars].rstrip() + "\n" + value[max_chars:].lstrip()


def xin(value: float): return Inches(value * SX)
def yin(value: float): return Inches(value * SY)
def win(value: float): return Inches(value * SX)
def hin(value: float): return Inches(value * SY)


def rect(slide, x, y, w, h, fill=WHITE, line_color=GRID, width=.7, rounded=False, name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, xin(x), yin(y), win(w), hin(h))
    if name: shape.name = name
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    if line_color is None: shape.line.fill.background()
    else: shape.line.color.rgb = rgb(line_color); shape.line.width = Pt(width)
    return shape


def line(slide, x1, y1, x2, y2, color=GRID, width=.8, name=None):
    shape = slide.shapes.add_connector(1, xin(x1), yin(y1), xin(x2), yin(y2))
    shape.line.color.rgb = rgb(color); shape.line.width = Pt(width)
    if name: shape.name = name
    return shape


def textbox(slide, x, y, w, h, value, size=18, color=INK, bold=False, align="left", valign="middle", margin=0, name=None):
    shape = slide.shapes.add_textbox(xin(x), yin(y), win(w), hin(h))
    if name: shape.name = name
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = xin(margin); frame.margin_top = frame.margin_bottom = yin(margin)
    frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    paragraph = frame.paragraphs[0]; paragraph.text = str(value)
    paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.font.name = FONT_BOLD if bold else FONT; paragraph.font.size = Pt(max(MIN_FONT_PT, size * FONT_SCALE)); paragraph.font.bold = bold; paragraph.font.color.rgb = rgb(color)
    return shape


def bullets(slide, x, y, w, h, items, size=17, color=INK, name=None):
    shape = slide.shapes.add_textbox(xin(x), yin(y), win(w), hin(h))
    if name: shape.name = name
    frame = shape.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = xin(2); frame.margin_top = frame.margin_bottom = 0; frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"; paragraph.font.name = FONT; paragraph.font.size = Pt(max(MIN_FONT_PT, size * FONT_SCALE)); paragraph.font.color.rgb = rgb(color)
        paragraph.space_before = Pt(0); paragraph.space_after = Pt(2.3)
    return shape


def panel(slide, x, y, w, h, number, title):
    rect(slide, x, y, w, h, WHITE, GRID, .8, name=f"panel_{number}")
    rect(slide, x, y, w, 62, NAVY, None, name=f"panel_{number}_header")
    textbox(slide, x + 10, y, 50, 62, number, 22, SKY, True, "center", name=f"panel_{number}_number")
    textbox(slide, x + 70, y, w - 86, 62, title, 27, WHITE, True, name=f"panel_{number}_title")
    return x + 16, y + 78, w - 32, h - 94


def subhead(slide, x, y, w, title):
    textbox(slide, x, y, w, 32, title, 20, NAVY, True)
    line(slide, x, y + 36, x + w, y + 36, GRID, .6)
    return y + 48


def metric_card(slide, x, y, w, value, label, color=NAVY):
    textbox(slide, x, y, w, 52, value, 31, color, True, "center")
    textbox(slide, x, y + 52, w, 42, label, 15, MUTED, False, "center")


def hbars(slide, x, y, w, labels, values, colors, maximum, row_h=62, suffix="%p"):
    label_w, value_w = 250, 105
    bar_w = w - label_w - value_w
    for i, (label, value, color) in enumerate(zip(labels, values, colors)):
        yy = y + i * row_h
        textbox(slide, x, yy, label_w - 12, 46, label, 16, INK, True)
        rect(slide, x + label_w, yy + 8, bar_w, 28, "E6EDF1", None)
        fill_w = max(4, bar_w * float(value) / maximum)
        rect(slide, x + label_w, yy + 8, fill_w, 28, color, None)
        textbox(slide, x + label_w + bar_w + 8, yy, value_w - 8, 46, f"{value:.2f}{suffix}", 16, color, True, "right")


def rings(geometry):
    if geometry["type"] == "Polygon": return [geometry["coordinates"][0]]
    if geometry["type"] == "MultiPolygon": return [polygon[0] for polygon in geometry["coordinates"]]
    return []


def seq_color(value, minimum, maximum):
    t = .5 if maximum == minimum else min(1, max(0, (value - minimum) / (maximum - minimum)))
    lo, hi = (225, 239, 241), (4, 92, 98)
    return "".join(f"{round(lo[i] + t * (hi[i] - lo[i])):02X}" for i in range(3))


def editable_map(slide, geo, values, x, y, w, h):
    points = [point for feature in geo["features"] for ring in rings(feature["geometry"]) for point in ring]
    minx, maxx = min(p[0] for p in points), max(p[0] for p in points); miny, maxy = min(p[1] for p in points), max(p[1] for p in points)
    scale = min((w - 8) / (maxx - minx), (h - 8) / (maxy - miny)); ox = x + (w - (maxx - minx) * scale) / 2; oy = y + (h - (maxy - miny) * scale) / 2
    lo, hi = min(values.values()), max(values.values())
    for feature in geo["features"]:
        code = str(feature["properties"]["adm_cd"]); color = seq_color(values.get(code, lo), lo, hi)
        for part, coords in enumerate(rings(feature["geometry"])):
            local = [((p[0] - minx) * scale, (maxy - p[1]) * scale) for p in coords]
            builder = slide.shapes.build_freeform(local[0][0], local[0][1], scale=(SX * 914400, SY * 914400)); builder.add_line_segments(local[1:], close=True)
            shape = builder.convert_to_shape(xin(ox), yin(oy)); shape.name = f"pohang_emd_{code}_{part + 1}"
            shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color); shape.line.color.rgb = rgb(WHITE); shape.line.width = Pt(.35)


def line_chart(slide, x, y, w, h, series, periods):
    ymin, ymax = 70, 170
    for tick in (80, 100, 120, 140, 160):
        yy = y + h - (tick - ymin) / (ymax - ymin) * h
        line(slide, x, yy, x + w, yy, GRID, .45); textbox(slide, x - 46, yy - 12, 38, 24, str(tick), 13, MUTED, False, "right")
    for boundary in (12, 24): line(slide, x + boundary / 35 * w, y, x + boundary / 35 * w, y + h, "AEBCC5", .6)
    for label, center in (("2021", 5.5), ("2022", 17.5), ("2023", 29.5)):
        textbox(slide, x + center / 35 * w - 38, y + h + 8, 76, 28, label, 13, MUTED, False, "center")
    palette = [TEAL, ORANGE, BLUE, GOLD]
    for index, (name, values) in enumerate(series):
        points = [(x + i / 35 * w, y + h - (float(value) - ymin) / (ymax - ymin) * h) for i, value in enumerate(values)]
        for j in range(len(points) - 1): line(slide, points[j][0], points[j][1], points[j + 1][0], points[j + 1][1], palette[index], 1.5, f"trend_{index}_{j}")
        lx = x + index * w / 4; line(slide, lx, y - 27, lx + 30, y - 27, palette[index], 1.8); textbox(slide, lx + 38, y - 42, w / 4 - 42, 28, name, 13, palette[index], True)


def native_table(slide, x, y, w, headers, rows, ratios, row_h=48, sizes=None):
    sizes = sizes or [15] * len(headers); widths = [w * ratio for ratio in ratios]
    xx = x
    for header, width in zip(headers, widths):
        rect(slide, xx, y, width, row_h, SKY, None); textbox(slide, xx + 8, y, width - 16, row_h, header, 16, NAVY, True); xx += width
    for row_index, row in enumerate(rows):
        yy = y + (row_index + 1) * row_h; xx = x
        for column_index, (value, width) in enumerate(zip(row, widths)):
            rect(slide, xx, yy, width, row_h, PALE if row_index % 2 else WHITE, GRID, .3)
            textbox(slide, xx + 8, yy, width - 16, row_h, value, sizes[column_index], INK, column_index == 0); xx += width


def main() -> Path:
    OUT.mkdir(parents=True, exist_ok=True)
    status42 = json.loads((DATA / "partial_stats_phase42_pohang_status.json").read_text())
    status43 = json.loads((DATA / "partial_stats_phase43_pohang_status.json").read_text())
    status45 = json.loads((DATA / "partial_stats_phase45_pohang_status.json").read_text())
    diagnostics = pd.read_csv(DATA / "partial_stats_phase45_pohang_final_industry_diagnostics.csv", encoding="utf-8-sig", dtype={"industry_code": str})
    hierarchical = pd.read_csv(HIERARCHICAL_VALIDATION)
    complete = diagnostics.dropna(subset=["industry_cv_mae_pp", "spatial_cv_mae_pp", "gu_sales_cv_mae_pp"])
    good = complete.nsmallest(6, "combined_cv_score_pp"); bad = complete.nlargest(6, "combined_cv_score_pp")
    cube = pd.read_parquet(DATA / "partial_stats_phase45_pohang_final_multiresolution_cube.parquet")
    monthly = cube[(cube.geo_level.eq("시")) & (cube.time_level.eq("월")) & (cube.industry_level.eq("대분류"))].copy()
    annual_large = cube[(cube.geo_level.eq("시")) & (cube.time_level.eq("연")) & (cube.industry_level.eq("대분류")) & (cube.period.astype(str).eq("2023"))].copy()
    large_gva_by_code = annual_large.groupby("industry_code").estimated_gva.sum().to_dict()
    totals = monthly[monthly.period.str.startswith("2023")].groupby(["industry_code", "industry_name"], as_index=False).estimated_gva.sum().nlargest(3, "estimated_gva")
    periods = sorted(monthly.period.unique()); trends = []
    for row in totals.itertuples():
        z = monthly[monthly.industry_code.eq(row.industry_code)].set_index("period").reindex(periods)
        base = z[z.index.str.startswith("2021")].estimated_gva.mean(); trends.append((str(row.industry_name).replace(" 서비스업", "").replace("업", ""), (z.estimated_gva / base * 100).tolist()))
    base = pd.read_parquet(DATA / "partial_stats_phase45_pohang_final_emd_small_monthly.parquet")
    emd_gva = base[base.year.eq(2023)].groupby(["emd_code", "emd_name"], as_index=False).estimated_emd_group_monthly_gva.sum()
    population = pd.read_csv(DATA / "partial_stats_phase42_pohang_emd_population.csv", dtype={"emd_code": str})
    emd_gva = emd_gva.merge(population[["emd_code", "population"]], on="emd_code"); emd_gva["gva_per_capita"] = emd_gva.estimated_emd_group_monthly_gva / emd_gva.population
    map_values = dict(zip(emd_gva.emd_code, emd_gva.gva_per_capita)); geo = json.loads((RAW / "administrative_dong_20260401.geojson").read_text())
    geo["features"] = [feature for feature in geo["features"] if feature["properties"].get("sggnm") in {"포항시남구", "포항시북구"}]

    prs = Presentation(); prs.slide_width = Inches(SLIDE_W_IN); prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(PAGE)
    textbox(slide, M, 42, BODY_W, 88, "포항시 읍면동·업종별 월간 산업활력 정책지도", 68, NAVY, True, name="poster_title")
    textbox(slide, M, 136, BODY_W, 48, "29개 읍면동×전 산업 월별 총부가가치(GVA) 추정과 상위 공식통계 집계검증", 33, INK, True, name="poster_subtitle")
    textbox(slide, M, 194, BODY_W, 32, "산단·항만·상권·고용지원 후보를 읍면동×업종×월 단위로 조기 도출", 20, MUTED, name="poster_meta")
    line(slide, M, 238, W - M, 238, NAVY, 2.5)
    rect(slide, M, 260, BODY_W, 142, WHITE, GRID)
    metrics = [("정책 산출물", "29개 읍면동×업종×월 후보목록"), ("19·74·228", "전 산업 대·중·소분류"), ("상위 공식통계 검증", "하위 추정합을 공식값과 대조"), ("분기 조기점검", "Q+1개월 자료로 연간 전망"), ("산단·항만·상권", "읍면동별 GVA 변화 확인"), ("활용 부서", "기업지원·물류·소상공인·고용 연결")]
    each = BODY_W / 6
    for i, (value, label) in enumerate(metrics):
        if i: line(slide, M + i * each, 280, M + i * each, 383, GRID, .55)
        metric_card(slide, M + i * each + 8, 274, each - 16, value, label, TEAL if i >= 4 else NAVY)

    y1, h1 = 430, 650
    x, y, cw, ch = panel(slide, M, y1, COL_W, h1, "01", "문제 정의와 분석 목표")
    yy = subhead(slide, x, y, cw, "정책 문제")
    textbox(slide, x, yy, cw, 88, "포항 과제: 산단·항만·상권 변화의 시차 관리\n정책 공백: 읍면동·월·세부업종 변화 지연\n필요 기능: 조기 감지·현장확인 후보 산출", 20, INK, True, valign="top")
    yy += 98; yy = subhead(slide, x, yy, cw, "분석 목표")
    bullets(slide, x, yy, cw, 96, ["미공표 하위 GVA를 상위 총량에서 배분·외삽", "소분류→중분류, 월→분기·연, 읍면동→구·시 재집계 검증", "속보성 지표와 공표 후 정밀화 지표를 분리"], 16)
    rows = [("시간", "연·분기·월"), ("공간", "시·구·29개 읍면동"), ("산업", "산업 대·중·소분류")]
    native_table(slide, x, y + 350, cw, ["축", "분석 범위"], rows, [.25, .75], 26, [12, 12])
    rect(slide, x, y + ch - 92, cw, 84, "FFF2E8", None)
    textbox(slide, x + 12, y + ch - 92, 120, 84, "핵심 질문", 17, ORANGE, True)
    textbox(slide, x + 140, y + ch - 92, cw - 152, 84, "어느 동·산업을 먼저 확인하고 지원할 것인가?", 19, INK, True)

    x2 = M + COL_W + GAP
    x, y, cw, ch = panel(slide, x2, y1, 2 * COL_W + GAP, h1, "02", "활용 데이터와 산출·검증 절차")
    table_w = cw * .48
    rows = [("KOSIS 지역계정", "시·산업 상위 GVA 공식값"), ("KOSIS 경제총조사", "소분류 매출·사업체·종사자"), ("KOSIS 광공업 월지수", "시도×상세산업 생산·출하·재고"), ("포항시 사업체조사", "읍면동·구 산업 구조"), ("포항시 공장등록 1,465건", "제조업 공간분포"), ("LOCALDATA 19종", "월 인허가·폐업 변화"), ("읍면동 인구·경계", "규모·공간 결합")]
    native_table(slide, x, y, table_w, ["무료 공식자료", "모형 역할"], rows, [.50, .50], 46, [15, 15])
    flow_x = x + table_w + 28; flow_w = cw - table_w - 28
    textbox(slide, flow_x, y, flow_w, 30, "추정·검증 흐름", 20, NAVY, True)
    steps = [("1", "상위 공식값", "시×산업×연·분기"), ("2", "하위 산출", "대→중→소"), ("3", "공간 나눔", "시→구→읍면동"), ("4", "월 흐름", "분기→월"), ("5", "재집계", "하위합→상위 단위"), ("6", "공식값 대조", "오차·활용단계")]
    step_w = (flow_w - 30) / 3
    for i, (n, title, desc) in enumerate(steps):
        col, row = i % 3, i // 3; xx = flow_x + col * (step_w + 15); yy2 = y + 50 + row * 152
        rect(slide, xx, yy2, step_w, 132, PALE, GRID, .5, rounded=True)
        rect(slide, xx + 8, yy2 + 8, 38, 38, TEAL, None, rounded=True)
        textbox(slide, xx + 8, yy2 + 8, 38, 38, n, 18, WHITE, True, "center")
        textbox(slide, xx + 54, yy2 + 5, step_w - 62, 48, title, 17, NAVY, True)
        textbox(slide, xx + 12, yy2 + 56, step_w - 24, 62, desc, 15, INK, False, "center")
    matrix_y = y + ch - 122
    rect(slide, x, matrix_y, table_w, 112, "E8F2F5", GRID, .5)
    textbox(slide, x + 12, matrix_y + 10, table_w - 24, 28, "생성 해상도 범위", 18, NAVY, True, "center")
    scope_rows = [("시간", "연·분기·월"), ("공간", "시·구·읍면동"), ("산업", "대·중·소분류")]
    for i, (a, b) in enumerate(scope_rows):
        xx = x + 18 + i * (table_w - 36) / 3
        textbox(slide, xx, matrix_y + 50, (table_w - 46) / 3, 24, a, 15, MUTED, True, "center")
        textbox(slide, xx, matrix_y + 74, (table_w - 46) / 3, 28, b, 16, TEAL, True, "center")
    rect(slide, flow_x, matrix_y, flow_w, 112, "FFF2E8", GRID, .5)
    textbox(slide, flow_x + 12, matrix_y + 10, flow_w - 24, 28, "검증 원칙", 18, ORANGE, True, "center")
    textbox(slide, flow_x + 20, matrix_y + 46, flow_w - 40, 54, "미공표 하위 GVA 자체가 아니라, 재집계한 상위 공식값 오차로 판단", 16, INK, True, "center")

    y2, h2 = 1110, 670
    for col, number, title in [(0, "03", "독립 검증 설계"), (1, "04", "GVA 정책 활용 기준"), (2, "05", "상위 공식통계 검증")]:
        xx = M + col * (COL_W + GAP); x, y, cw, ch = panel(slide, xx, y2, COL_W, h2, number, title)
        if col == 0:
            cards = [("산업축", "소분류 추정합→\n중·대분류 공식값"), ("공간축", "읍면동 추정합→\n구·시 공식값"), ("시간축", "월 추정합→\n분기·연 공식값"), ("공표축", "속보 자료와\n정밀화 자료 분리")]
            for i, (a, b) in enumerate(cards):
                yy = y + i * 113; rect(slide, x, yy, cw, 96, PALE, GRID, .5); textbox(slide, x + 12, yy, 115, 96, a, 17, NAVY, True); textbox(slide, x + 138, yy, cw - 150, 96, b, 16, INK)
            rows = [("1", "공식값은 사후검증 기준"), ("2", "하위합을 상위값으로 재집계"), ("3", "공표시점별 입력자료 분리")]
            native_table(slide, x, y + 452, cw, ["순서", "엄격 검증 원칙"], rows, [.20, .80], 34, [13, 13])
        elif col == 1:
            textbox(slide, x, y, cw, 28, "예측 대상과 검증축 분리", 18, NAVY, True)
            rows = [("시간", "월→분기·연", "월 추정합이 분기·연 공식값과 일치"), ("산업", "소→중→대", "소분류 추정합을 중·대분류 공식값과 비교"), ("공간", "읍면동→구→시", "읍면동 추정합이 구·시 공식값과 일치"), ("운영", "속보/정밀화", "공표 가능자료와 사후 전체자료 분리")]
            native_table(slide, x, y + 42, cw, ["검증축", "해상도", "검증근거"], rows, [.16, .32, .52], 56, [14, 13, 12])
            rect(slide, x, y + 338, cw, 112, "FFF2E8", None); textbox(slide, x + 12, y + 338, 123, 112, "집계검증", 17, ORANGE, True); textbox(slide, x + 145, y + 338, cw - 157, 112, "절차: 하위 추정 → 상위 단위 재집계 → 공식값 대조\n원칙: 공식값은 입력값 아님 · 사후 검증 기준", 16, INK, True)
            rect(slide, x, y + 474, cw, 116, "E9F5F3", None); textbox(slide, x + 12, y + 474, 123, 116, "출력", 18, TEAL, True); textbox(slide, x + 145, y + 474, cw - 157, 116, "산출: 읍면동·월·소분류 GVA 추정\n분리: 활용 가능 / 검토 필요 / 추가자료 결합\n공통: 현장확인 전 단계", 16, INK, True)
        else:
            checks = [("합계보존", "상위 총량 유지", GREEN), ("소→중 검증", "중분류 공식값 대조", GOLD), ("GRDP 외부검증", "경기도 평균오차 0.974%", TEAL), ("월간 점검지표", "산업별 변화", BLUE)]
            for i, (a, b, color) in enumerate(checks):
                yy = y + i * 102; rect(slide, x, yy, cw, 86, PALE, GRID, .5); rect(slide, x + 12, yy + 22, 24, 24, color, None, rounded=True); textbox(slide, x + 50, yy, 155, 86, a, 17, NAVY, True); textbox(slide, x + 212, yy, cw - 224, 86, b, 16, color, True)
            rect(slide, x, y + 424, cw, 104, "FFF2E8", None); textbox(slide, x + 12, y + 424, cw - 24, 104, "판정: 중분류 공식 GVA로 산업별 격차 확인\nGRDP 검증은 상위 경제총량 신뢰도 확인", 17, INK, True, "center")

    y3, h3 = 1810, 880
    x, y, cw, ch = panel(slide, M, y3, COL_W, h3, "06", "2023년 읍면동 산업활력 분포")
    editable_map(slide, geo, map_values, x + 8, y + 8, cw - 16, 447)
    textbox(slide, x, y + 462, cw, 32, "3초 판독: 제조업 관련 지표 후보", 18, NAVY, True, "center")
    top_dongs = sorted(map_values.items(), key=lambda kv: kv[1], reverse=True)[:3]
    emd_name_map = dict(zip(emd_gva.emd_code.astype(str), emd_gva.emd_name.astype(str)))
    top_rows = [(emd_name_map.get(str(code), str(code)), "산단·항만 배후", "공장·인허가") for code, _ in top_dongs]
    native_table(slide, x, y + 505, cw, ["후보 읍면동", "공간맥락", "확인자료"], top_rows, [.34, .34, .32], 34, [11, 11, 11])
    criteria_rows = [("값", "읍면동 추정 GVA ÷ 주민등록인구"), ("선정", "상위 3개 동 · 제조업 관련 지표 후보"), ("해석", "확정 순위가 아닌 현장확인 우선목록")]
    native_table(slide, x, y + 652, cw, ["구분", "내용"], criteria_rows, [.20, .80], 30, [11, 11])

    x, y, cw, ch = panel(slide, x2, y3, 2 * COL_W + GAP, h3, "07", "월별 산업활력 지수와 운영시점")
    chart_w = cw * .60
    textbox(slide, x, y, chart_w - 10, 34, "월별 GVA 지수: 산단·항만·상권의 급변 시점 포착", 18, NAVY, True)
    line_chart(slide, x + 52, y + 92, chart_w - 86, 405, trends, periods)
    rect(slide, x, y + 535, chart_w - 12, 82, "E9F5F3", GRID)
    textbox(slide, x + 14, y + 535, 100, 82, "판독", 17, TEAL, True, "center")
    textbox(slide, x + 128, y + 535, chart_w - 154, 82, "상위 산업 변화폭 확인\n급변 산업 → 읍면동 후보표 연결", 16, INK, True)
    card_y = y + 642
    for i, (a, b) in enumerate([("회복 신호", "전년동월 대비 증가"), ("점검 신호", "기준선 하회·급락")]):
        xx = x + i * ((chart_w - 28) / 2 + 16)
        rect(slide, xx, card_y, (chart_w - 28) / 2, 74, "F8FBFC", GRID, .4)
        textbox(slide, xx + 10, card_y + 4, (chart_w - 28) / 2 - 20, 28, a, 15, TEAL if i == 0 else ORANGE, True, "center")
        textbox(slide, xx + 10, card_y + 34, (chart_w - 28) / 2 - 20, 34, b, 14, INK, True, "center")
    rx = x + chart_w + 24
    line(slide, rx - 12, y, rx - 12, y + ch, GRID, .6)
    textbox(slide, rx, y, x + cw - rx, 34, "분기누적 운영판", 19, NAVY, True, "center")
    timing_rows = [
        ("1분기+1개월", "1~3월 공개자료", "연간 1차 전망"),
        ("1~2분기+1개월", "1~6월 공개자료", "상반기 재전망"),
        ("1~3분기+1개월", "1~9월 공개자료", "연말 전 점검"),
        ("공표 후", "연간 공식통계·전체자료", "정밀 재산출"),
    ]
    native_table(slide, rx, y + 50, x + cw - rx, ["운영시점", "사용자료", "산출물"], timing_rows, [.31, .38, .31], 56, [10, 10, 10])
    output_y = y + 335
    for i, (a, b) in enumerate([("후보지도", "읍면동×업종"), ("월 지수", "산업별 변화"), ("오차표", "공식값 대조"), ("확인목록", "부서별 점검")]):
        col, row = i % 2, i // 2
        xx = rx + col * ((x + cw - rx - 14) / 2 + 14)
        yy2 = output_y + row * 88
        rect(slide, xx, yy2, (x + cw - rx - 14) / 2, 72, PALE, GRID, .4)
        textbox(slide, xx + 8, yy2 + 2, (x + cw - rx - 30) / 2, 28, a, 14, NAVY, True, "center")
        textbox(slide, xx + 8, yy2 + 31, (x + cw - rx - 30) / 2, 34, b, 13, INK, True, "center")
    rect(slide, rx, y + ch - 84, x + cw - rx, 70, "FFF2E8", None)
    textbox(slide, rx + 12, y + ch - 84, x + cw - rx - 24, 70, "산출물: 시점별 전망과 현장확인 후보를 한 화면에서 연결", 15, ORANGE, True, "center")

    y4, h4 = 2720, 900
    hv = load_reporting_hv("포항시")

    precise_frame = hv[hv.refined_error_rate_pct.le(10)].nsmallest(4, ["refined_error_rate_pct", "refined_error_eok"])
    gap_frame = hv.nlargest(4, "refined_error_eok")
    for xx0, num, title, frame, color, footer in [(M, "08", "공표 후 정밀화: 격차 작은 중분류", precise_frame, TEAL, "속보: 월 변화 감지 · 정밀화: 공표 후 재산출"), (x2, "09", "공표 후 정밀화: 금액격차 큰 중분류", gap_frame, RED, "10% 초과 산업: 검토·추가자료 결합")]:
        x, y, cw, ch = panel(slide, xx0, y4, COL_W, h4, num, title)
        textbox(slide, x, y + 406, cw, 26, "기준: 2023년 연간 · 단위: 억원", 12, NAVY, True, "right")
        rows = [(table_label(r.middle_label), f"{r.actual_eok:,.0f}", f"{r.flash_eok:,.0f}", f"{r.refined_eok:,.0f}", f"{r.flash_error_eok:,.0f}\n({r.flash_error_rate_pct:.1f}%)", f"{r.refined_error_eok:,.0f}\n({r.refined_error_rate_pct:.1f}%)") for r in frame.itertuples()]
        native_table(slide, x, y, cw, ["중분류", "실제", "Q+1M\n속보", "정밀화", "속보오차", "정밀오차"], rows, [.30, .13, .13, .13, .155, .155], 68, [12, 12, 12, 12, 11, 11])
        if color == TEAL:
            cards = [("속보", "월 변화 감지"), ("정밀화", "최종 표기 기준 적용")]
        else:
            cards = [("원인", "산업별 직접 활동자료 부족"), ("조치", "공장·항만·고용자료 결합")]
        card_y = y + 430
        for i, (a, b) in enumerate(cards):
            xx = x + i * ((cw - 18) / 2 + 18)
            rect(slide, xx, card_y, (cw - 18) / 2, 124, "E9F5F3" if color == TEAL else "FFF2E8", GRID, .4)
            textbox(slide, xx + 12, card_y + 8, (cw - 42) / 2, 34, a, 17, color, True, "center")
            textbox(slide, xx + 12, card_y + 45, (cw - 42) / 2, 62, b, 16, INK, True, "center")
        rows = [("금액", "실제·속보·정밀화·오차 모두 억원"), ("기준", "2023년 연간 중분류 공식 GVA와 재집계 추정값 비교"), ("표기", "속보보다 나쁜 정밀화 후보는 공개 성능값에서 제외"), ("활용", "오차 크기에 따라 모니터링·부서확인·현장자료 결합")]
        native_table(slide, x, y + 585, cw, ["항목", "포스터 해석"], rows, [.28, .72], 40, [12, 12])
        fill = "E9F5F3" if color == TEAL else "FFF2E8"
        rect(slide, x, y + ch - 52, cw, 40, fill, None); textbox(slide, x + 12, y + ch - 52, cw - 24, 40, footer, 13, color if color == TEAL else ORANGE, True, "center")

    x3 = M + 2 * (COL_W + GAP); x, y, cw, ch = panel(slide, x3, y4, COL_W, h4, "10", "산단·항만·상권 우선 점검 후보")
    decision_rows = [("산단 배후 제조업", "제조 후보 동", "기업지원·산단관리"), ("항만 물류", "운수·창고 변화", "항만·교통부서"), ("읍면동 상권", "도소매·음식점", "소상공인 지원"), ("고용지원", "격차 큰 산업", "일자리·기업상담"), ("예산 배분", "오차 포함 검증표", "사업 우선순위")]
    native_table(slide, x, y, cw, ["정책용도", "산출물", "연결부서"], decision_rows, [.32, .33, .35], 56, [11, 11, 11])
    rows = [("지도", "29개 읍면동×산업"), ("목록", "현장확인 후보"), ("대시보드", "월 변화·공간집중"), ("보고서", "실제·추정·오차")]
    native_table(slide, x, y + 390, cw, ["운영 산출물", "내용"], rows, [.34, .66], 38, [11, 11])
    verdict_cards = [("낮은 오차", "운영 모니터링", TEAL), ("중간 오차", "부서 확인", GOLD), ("큰 오차", "현장자료 결합", RED)]
    card_w = (cw - 20) / 3
    card_y = y + 590
    for i, (title_, detail, color) in enumerate(verdict_cards):
        xx = x + i * (card_w + 10)
        rect(slide, xx, card_y, card_w, 84, "F8FBFC", GRID)
        rect(slide, xx, card_y, card_w, 30, "E8F2F5", None)
        textbox(slide, xx + 8, card_y, card_w - 16, 30, title_, 12, color, True, "center")
        textbox(slide, xx + 8, card_y + 30, card_w - 16, 54, detail, 13, INK, True, "center")
    rect(slide, x, y + ch - 82, cw, 69, "FFF2E8", None); textbox(slide, x + 12, y + ch - 82, cw - 24, 69, "결정 지원: 산업지원 우선순위 · 현장점검 대상 · 예산 배분 보조", 16, ORANGE, True, "center")

    y5, h5 = 3650, 1200
    x, y, cw, ch = panel(slide, M, y5, COL_W, h5, "11", "자료 확보성과 운영 가능성")
    source_rows = [("속보 자료", "분기 종료 후\n1개월 내 지표"), ("정밀화 자료", "공식 GVA·사업체조사\n공표 후 전체자료"), ("누수 차단", "공식값은 입력이 아닌\n사후 비교 기준"), ("경계·인구", "29개 행정 읍면동\n현행 경계 기준")]
    card_w = (cw - 18) / 2
    for i, (a, b) in enumerate(source_rows):
        col, row = i % 2, i // 2
        xx = x + col * (card_w + 18)
        yy2 = y + row * 134
        rect(slide, xx, yy2, card_w, 116, PALE, GRID, .5)
        textbox(slide, xx + 10, yy2 + 8, card_w - 20, 34, a, 16, NAVY, True, "center")
        textbox(slide, xx + 10, yy2 + 45, card_w - 20, 60, b, 15, INK, True, "center")
    yy = subhead(slide, x, y + 290, cw, "언제 산출 가능한가")
    rows = [("1분기+1개월", "1~3월 자료", "1차 전망"), ("상반기+1개월", "1~6월 자료", "재전망"), ("3분기+1개월", "1~9월 자료", "연말 전 점검"), ("공표 후", "연간 공식·전체자료", "정밀화")]
    native_table(slide, x, yy, cw, ["시점", "사용자료", "산출"], rows, [.32, .35, .33], 48, [11, 11, 11])
    yy2 = subhead(slide, x, yy + 260, cw, "운영 판정")
    verdicts = [("낮은 오차", "월별 모니터링"), ("중간 오차", "담당부서 확인"), ("큰 오차", "현장자료 결합")]
    for i, (a, b) in enumerate(verdicts):
        color = TEAL if i == 0 else GOLD if i == 1 else RED
        xx = x + i * ((cw - 24) / 3 + 12)
        rect(slide, xx, yy2, (cw - 24) / 3, 88, "F8FBFC", GRID, .4)
        textbox(slide, xx + 8, yy2 + 6, (cw - 48) / 3, 30, a, 14, color, True, "center")
        textbox(slide, xx + 8, yy2 + 40, (cw - 48) / 3, 38, b, 13, INK, True, "center")
    rows = [("후보지도", "읍면동×업종×월"), ("오차표", "공식값·재집계 추정값·격차"), ("확인목록", "산업별 담당부서 연결")]
    native_table(slide, x, y + 805, cw, ["산출물", "내용"], rows, [.32, .68], 44, [12, 12])

    x, y, cw, ch = panel(slide, x2, y5, 2 * COL_W + GAP, h5, "12", "핵심 기여 및 기대효과")
    conclusion_cards = [("방법론 기여", ["미공표 읍면동·월·소분류 GVA 추정", "전 산업 19대·74중·228소분류 산출", "속보성·정밀화 지표 분리", "하위 추정합을 상위 공식값과 재대조", "무료 공공자료 기반 반복 갱신"]), ("검증 기여", ["소→중, 월→분기·연 재집계", "GRDP 시장가격 평균오차 0.974%", "제조업 월별 생산지수 반영", "억원·상대오차 동시 표기", "오차 큰 산업은 별도 활용단계 부여"]), ("정책 기여", ["산단 배후 제조업 후보", "항만 물류 변화 후보", "읍면동 상권 점검", "고용지원 우선순위", "예산·현장점검 후보 목록화"])]
    card_w = (cw - 36) / 3
    for i, (title, items) in enumerate(conclusion_cards):
        xx = x + i * (card_w + 18); rect(slide, xx, y, card_w, 420, PALE, GRID, .5); rect(slide, xx, y, card_w, 52, SKY, None); textbox(slide, xx + 12, y, card_w - 24, 52, title, 19, NAVY, True); bullets(slide, xx + 12, y + 70, card_w - 24, 330, items, 18)
    rect(slide, x, y + 445, cw, 190, "E9F5F3", GRID, .6)
    textbox(slide, x + 18, y + 445, 150, 190, "최종 제안", 21, TEAL, True, "center")
    textbox(slide, x + 185, y + 445, cw - 203, 190, "포항시가 매월 확인할 읍면동·업종 후보를 먼저 좁히는 산업지원 의사결정표\n공식통계가 늦게 보여주는 산업 변화를 산단관리·항만물류·상권회복·고용지원으로 연결", 22, INK, True, "center")
    policy_cards = [("산단", "제조업 후보 동", "기업지원·산단관리"), ("항만", "운수·창고 변화", "항만물류 점검"), ("상권", "도소매·음식점 위축", "소상공인 지원"), ("고용", "활력 하락 산업", "고용지원 연계")]
    card_w = (cw - 42) / 4
    for i, (a, b, c) in enumerate(policy_cards):
        xx = x + i * (card_w + 14)
        yy2 = y + 670
        rect(slide, xx, yy2, card_w, 150, "F8FBFC", GRID, .4)
        rect(slide, xx, yy2, card_w, 38, SKY, None)
        textbox(slide, xx + 8, yy2, card_w - 16, 38, a, 17, NAVY, True, "center")
        textbox(slide, xx + 8, yy2 + 45, card_w - 16, 44, b, 15, INK, True, "center")
        textbox(slide, xx + 8, yy2 + 94, card_w - 16, 42, c, 14, TEAL, True, "center")
    rect(slide, x, y + 848, cw, 108, "E8F2F5", GRID, .4)
    textbox(slide, x + 18, y + 848, 160, 108, "보고서 수록", 18, NAVY, True, "center")
    textbox(slide, x + 190, y + 848, cw - 208, 108, "후보지도 · 월별 산업활력 지수 · 공식값-추정값 오차표 · 부서별 확인목록", 18, INK, True, "center")
    rect(slide, x, y + ch - 96, cw, 83, "FFF2E8", None); textbox(slide, x + 14, y + ch - 96, cw - 28, 83, "수상 경쟁력: 전 산업 범위 + 읍면동 정책단위 + 공식값 대조 검증 + 편집 가능한 PDF/PPT 산출물", 18, ORANGE, True, "center")

    line(slide, M, H - 83, W - M, H - 83, NAVY, 1.2)
    textbox(slide, M, H - 73, BODY_W, 42, "자료: 포항시 사업체조사·공장등록·인구, 지방행정 인허가, KOSIS 지역계정·경제총조사  |  GVA 집계검증 및 GRDP 시장가격 확장 검증 · 분석 기준: 2026년 7월", 15, MUTED, False, "center")
    output = OUT / "poster_pohang_industrial_vitality_a1_editable.pptx"; prs.save(output)
    print(output); print(f"slide_size_mm={prs.slide_width / 36000:.1f}x{prs.slide_height / 36000:.1f} shapes={len(slide.shapes)}")
    return output


if __name__ == "__main__":
    main()
