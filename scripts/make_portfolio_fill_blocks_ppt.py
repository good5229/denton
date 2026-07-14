from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from make_portfolio_ppt import (
    BG,
    BLUE,
    BORDER,
    GOLD,
    GREEN,
    INK,
    LIGHT,
    MUTED,
    PURPLE,
    RED,
    SLATE,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    WHITE,
    add_circle,
    add_rect,
    add_text,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "ppt"


def add_feature_card(slide, x, y, w, h, num, title, body, accent) -> None:
    add_rect(slide, x, y, w, h, WHITE, BORDER)
    add_circle(slide, x + 24, y + 22, 42, accent, num)
    add_text(slide, x + 82, y + 22, w - 105, 32, title, 15, accent, True)
    add_text(slide, x + 82, y + 62, w - 105, 30, body, 11, MUTED, True)


def add_bottom_card(slide, x, y, w, h, title, value, note, accent) -> None:
    add_rect(slide, x, y, w, h, WHITE, BORDER)
    add_text(slide, x + 28, y + 18, 180, 30, title, 15, accent, True)
    add_text(slide, x + 28, y + 58, 180, 42, value, 20, INK, True)
    add_text(slide, x + 28, y + 108, w - 55, 28, note, 11, MUTED, True)


def build_copy_ready_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1920, 1080, BG, BG, False, 0)
    add_text(slide, 80, 55, 900, 62, "포트폴리오 보조 블록", 28, INK, True)
    add_text(slide, 82, 122, 1200, 36, "우측 빈 공간과 하단 빈 공간에 복사해 넣을 수 있는 개별 편집 요소", 15, MUTED, True)

    add_text(slide, 80, 205, 420, 40, "우측 패널 후보", 21, INK, True)
    add_rect(slide, 80, 260, 520, 580, "F9FBFD", BORDER)
    add_text(slide, 118, 292, 300, 34, "구현 난점과 해결", 18, INK, True)
    cards = [
        ("1", "이종 통계 결합", "GVA·지표 표준 키 정렬", BLUE),
        ("2", "지표 우선순위", "생산·서비스·건설·GDP", TEAL),
        ("3", "연간합 보존", "분기 합계 = 연간 기준", PURPLE),
        ("4", "외삽 검증", "실제값 vs 예측값", RED),
    ]
    y = 345
    for card in cards:
        add_feature_card(slide, 118, y, 445, 78, *card)
        y += 96

    add_text(slide, 720, 205, 460, 40, "하단 띠 후보", 21, INK, True)
    add_rect(slide, 720, 260, 1080, 310, "F9FBFD", BORDER)
    bottom = [
        ("데이터", "5년", "2019~2023", BLUE),
        ("범위", "전 업종", "지역×업종×분기", TEAL),
        ("방법", "Denton", "연간합+평활화", PURPLE),
        ("검증", "7.33%", "2023 MAPE", GOLD),
    ]
    x = 760
    for item in bottom:
        add_bottom_card(slide, x, 320, 235, 168, *item)
        x += 255

    add_text(slide, 720, 645, 460, 40, "짧은 설명문 후보", 21, INK, True)
    add_rect(slide, 720, 700, 1080, 140, WHITE, BORDER)
    add_text(slide, 760, 728, 980, 34, "연간 지역 총부가가치를 분기 지표 흐름에 맞춰 분기화했습니다.", 15, INK, True)
    add_text(slide, 760, 782, 980, 30, "핵심은 연간합 제약과 Denton 평활화를 동시에 만족하는 최적화입니다.", 13, MUTED, True)

    add_text(slide, 80, 910, 460, 34, "사용 팁", 18, SLATE, True)
    add_text(slide, 80, 950, 1260, 30, "각 카드·텍스트·원형 번호는 모두 분리된 PPT 요소입니다. 기존 슬라이드의 빈 공간 크기에 맞춰 그룹으로 복사하거나 개별 수정하면 됩니다.", 13, MUTED, True)


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_copy_ready_slide(prs)
    out = OUT / "portfolio_fill_blocks_editable.pptx"
    prs.save(out)
    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
