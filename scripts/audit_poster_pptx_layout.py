from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def text_overlaps(path: Path, threshold: float = 0.05) -> list[dict[str, object]]:
    presentation = Presentation(path)
    findings: list[dict[str, object]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        items = []
        for index, shape in enumerate(slide.shapes):
            value = shape.text.strip() if hasattr(shape, "text") else ""
            if not value:
                continue
            items.append(
                {
                    "index": index,
                    "name": shape.name,
                    "text": value.replace("\n", " / "),
                    "box": (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height),
                }
            )
        for left_index, left in enumerate(items):
            for right in items[left_index + 1 :]:
                x0 = max(left["box"][0], right["box"][0])
                y0 = max(left["box"][1], right["box"][1])
                x1 = min(left["box"][2], right["box"][2])
                y1 = min(left["box"][3], right["box"][3])
                if x1 <= x0 or y1 <= y0:
                    continue
                overlap_area = (x1 - x0) * (y1 - y0)
                left_area = (left["box"][2] - left["box"][0]) * (left["box"][3] - left["box"][1])
                right_area = (right["box"][2] - right["box"][0]) * (right["box"][3] - right["box"][1])
                ratio = overlap_area / min(left_area, right_area)
                if ratio > threshold:
                    findings.append(
                        {
                            "slide": slide_number,
                            "ratio": ratio,
                            "left": f"{left['name']}: {left['text'][:100]}",
                            "right": f"{right['name']}: {right['text'][:100]}",
                        }
                    )
    return findings


def main() -> None:
    parser = argparse.ArgumentParser(description="Fail when PPT text boxes geometrically overlap.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--threshold", type=float, default=0.05)
    args = parser.parse_args()
    findings = text_overlaps(args.pptx, args.threshold)
    if findings:
        for finding in sorted(findings, key=lambda row: float(row["ratio"]), reverse=True):
            print(
                f"slide={finding['slide']} ratio={finding['ratio']:.3f}\n"
                f"  {finding['left']}\n  {finding['right']}"
            )
        raise SystemExit(f"layout audit failed: {len(findings)} text overlaps")
    print("layout audit: PASS (0 text-box overlaps)")


if __name__ == "__main__":
    main()
