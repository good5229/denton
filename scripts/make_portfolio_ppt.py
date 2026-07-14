from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "ppt"

W, H = 1920, 1080
SLIDE_W, SLIDE_H = 13.333333, 7.5

BG = "F7F8FA"
INK = "17212B"
MUTED = "475569"
BORDER = "CAD2DE"
WHITE = "FFFFFF"
BLUE = "2F6FBB"
TEAL = "21867A"
GREEN = "50884F"
GOLD = "B07822"
RED = "B6544A"
PURPLE = "7A5FA8"
SLATE = "475569"
LIGHT = "EEF3F8"

FONT = "Apple SD Gothic Neo"


def emu_px(value: float) -> int:
    return Inches(value * SLIDE_W / W)


def color(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_run_font(run, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            rpr.append(el)
        el.set("typeface", font_name)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    fill: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(emu_px(x), emu_px(y), emu_px(w), emu_px(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = False
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, FONT)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = BORDER,
    radius: bool = True,
    line_width: float = 1.25,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, emu_px(x), emu_px(y), emu_px(w), emu_px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, fill: str, label: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu_px(x), emu_px(y), emu_px(d), emu_px(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(fill)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_run_font(run, FONT)
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color(WHITE)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    head = 16
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu_px(x1), emu_px(y1), emu_px(x2 - head), emu_px(y2))
    conn.line.color.rgb = color("6B7785")
    conn.line.width = Pt(2.5)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, emu_px(x2 - head), emu_px(y2 - head / 2), emu_px(head), emu_px(head))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = color("6B7785")
    tri.line.color.rgb = color("6B7785")


def add_divider(slide, x1: float, y: float, x2: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu_px(x1), emu_px(y), emu_px(x2), emu_px(y))
    line.line.color.rgb = color("D5DBE5")
    line.line.width = Pt(2)


def add_cell_bar(slide, x: float, y: float, w: float, h: float, labels: list[str], weights: list[float], fill: str) -> None:
    total = sum(weights)
    used = 0.0
    for idx, (label, weight) in enumerate(zip(labels, weights)):
        cw = w * weight / total if idx < len(labels) - 1 else w - used
        add_rect(slide, x + used, y, cw, h, fill, "145B55" if fill == TEAL else "173A63", False, 1.5)
        add_text(slide, x + used, y + 14, cw, h - 10, label, 15, WHITE, True, PP_ALIGN.CENTER)
        used += cw


def build_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, BG, BG, False, 0)

    add_text(slide, 84, 52, 1280, 74, "지역 GVA 분기화 엔진", 28, INK, True)
    add_rect(slide, 1460, 58, 285, 56, SLATE, SLATE)
    add_text(slide, 1460, 73, 285, 40, "Portfolio Tech", 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 90, 132, 1500, 48, "KOSIS 원천통계와 한국은행 방법론을 연결해 연간 지역 총부가가치를 분기 시계열로 복원", 18, MUTED, True)

    steps = [
        ("1", "수집", "KOSIS OpenAPI", "GVA·생산지수·GDP", BLUE),
        ("2", "매핑", "지역 × 업종 × 분기", "지표 우선순위 적용", TEAL),
        ("3", "추정", "비례형 덴튼", "연간합 제약 + 평활화", PURPLE),
        ("4", "검증", "2023 외삽 오차", "실제 GVA와 비교", RED),
    ]
    x = 88
    for idx, (num, title, line1, line2, accent) in enumerate(steps):
        add_rect(slide, x, 205, 380, 180)
        add_circle(slide, x + 26, 232, 56, accent, num)
        add_text(slide, x + 108, 225, 190, 44, title, 22, accent, True)
        add_text(slide, x + 108, 282, 250, 36, line1, 16, INK)
        add_text(slide, x + 108, 326, 260, 32, line2, 14, MUTED, True)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 390, 295, x + 438, 295)
        x += 455

    add_rect(slide, 110, 450, 1080, 410)
    add_text(slide, 155, 482, 220, 50, "핵심 구현", 23, INK, True)
    add_text(slide, 155, 543, 150, 40, "입력 벡터", 17, BLUE, True)
    add_text(slide, 360, 543, 230, 40, "Ay: 연간 GVA", 17, INK, True)
    add_cell_bar(slide, 610, 526, 410, 62, ["I₁", "I₂", "I₃", "I₄"], [22, 18, 29, 31], "4F84BF")
    add_text(slide, 1040, 545, 120, 30, "분기 지표", 14, MUTED, True)
    add_divider(slide, 170, 630, 1135)
    add_text(slide, 155, 674, 130, 36, "최적화", 17, PURPLE, True)
    add_text(slide, 360, 670, 640, 42, "min Σ [(Xt/It) - (Xt-1/It-1)]²", 17, INK, True)
    add_text(slide, 360, 728, 520, 38, "s.t. 각 연도 Σ분기 Xyq = Ay", 17, MUTED, True)
    add_text(slide, 155, 790, 130, 36, "산출", 17, TEAL, True)
    add_cell_bar(slide, 360, 772, 660, 62, ["X₁", "X₂", "X₃", "X₄"], [24, 19, 27, 30], TEAL)
    add_text(slide, 1040, 790, 120, 30, "분기 GVA", 14, MUTED, True)

    add_rect(slide, 1260, 450, 575, 410)
    add_text(slide, 1305, 482, 220, 50, "구현 결과", 23, INK, True)
    metrics = [("5,760", "분기 GVA 추정 행", GREEN), ("288", "2023 외삽 검증 건", RED), ("7.33%", "외삽 MAPE", GOLD)]
    y = 546
    for value, label, accent in metrics:
        add_text(slide, 1305, y, 190, 62, value, 33, accent, True)
        add_text(slide, 1518, y + 20, 260, 34, label, 17, MUTED, True)
        y += 92
    add_rect(slide, 1305, 805, 480, 36, LIGHT, LIGHT)
    add_text(slide, 1325, 811, 440, 30, "Python 파이프라인 + CSV", 13, MUTED)

    notes = [("데이터 처리", "API 수집 · 정규화"), ("알고리즘 구현", "Denton 최적화"), ("검증 설계", "외삽 오차 비교")]
    x = 135
    for title, body in notes:
        add_rect(slide, x, 910, 500, 78)
        add_text(slide, x + 28, 925, 230, 34, title, 16, BLUE, True)
        add_text(slide, x + 265, 928, 200, 30, body, 14, MUTED, True)
        x += 560


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_slide(prs)
    out = OUT / "portfolio_implementation_editable.pptx"
    prs.save(out)
    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
