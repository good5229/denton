from __future__ import annotations

import subprocess
from pathlib import Path

import pandas as pd
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "pohang" / "poster_pohang_industrial_vitality_a1_editable.pptx"
DATA = ROOT / "data" / "processed"


def main() -> int:
    presentation = Presentation(PPTX)
    assert len(presentation.slides) == 1
    assert abs(presentation.slide_width / 36000 - 594) < .1
    assert abs(presentation.slide_height / 36000 - 841) < .1
    slide = presentation.slides[0]
    text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
    for forbidden in ("Phase", "phase", "F00", "Q00", "C00", "nan", "NaN"):
        assert forbidden not in text, forbidden
    for required in (
        "포항시 산업활력 정밀지도", "29개 행정 읍면동", "KSIC 실제 업종명",
        "예측 양호 산업", "예측 취약 산업", "활용 판정 및 검증", "자료 확보성 검토",
        "연·분기·월", "시·구·읍면동", "대·중·소분류",
        "제조업", "창작 예술 및 여가관련 서비스업", "종합 건설업",
    ):
        assert required in text, required
    assert len(slide.shapes) >= 540
    cube = pd.read_parquet(DATA / "partial_stats_phase45_pohang_final_multiresolution_cube.parquet")
    expected = {
        ("시", "연", "대분류"), ("시", "연", "중분류"), ("시", "연", "소분류"),
        ("시", "분기", "대분류"), ("시", "분기", "중분류"), ("시", "분기", "소분류"),
        ("시", "월", "대분류"), ("시", "월", "중분류"), ("시", "월", "소분류"),
        ("구", "연", "대분류"), ("구", "연", "중분류"), ("구", "연", "소분류"),
        ("구", "분기", "대분류"), ("구", "분기", "중분류"), ("구", "분기", "소분류"),
        ("구", "월", "대분류"), ("구", "월", "중분류"), ("구", "월", "소분류"),
        ("읍면동", "연", "대분류"), ("읍면동", "연", "중분류"), ("읍면동", "연", "소분류"),
        ("읍면동", "분기", "대분류"), ("읍면동", "분기", "중분류"), ("읍면동", "분기", "소분류"),
        ("읍면동", "월", "대분류"), ("읍면동", "월", "중분류"), ("읍면동", "월", "소분류"),
    }
    actual = set(map(tuple, cube[["geo_level", "time_level", "industry_level"]].drop_duplicates().to_numpy()))
    assert actual == expected, actual ^ expected
    subprocess.run([str(ROOT / ".venv" / "bin" / "python"), str(ROOT / "scripts" / "audit_poster_pptx_layout.py"), str(PPTX)], check=True)
    print(f"Pohang poster verification: PASS size=594x841mm shapes={len(slide.shapes)} text_chars={len(text)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
