from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "ppt"

W, H = 1920, 1080
SLIDE_W, SLIDE_H = 13.333333, 7.5

BG = "F7F8FA"
INK = "17212B"
MUTED = "475569"
LINE = "CAD2DE"
WHITE = "FFFFFF"
BLUE = "2F6FBB"
TEAL = "21867A"
GREEN = "50884F"
GOLD = "B07822"
RED = "B6544A"
PURPLE = "7A5FA8"
LIGHT = "EEF3F8"
FONT = "Apple SD Gothic Neo"


def emu_px(value: float) -> int:
    return Inches(value * SLIDE_W / W)


def color(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_run_font(run, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            rpr.append(el)
        el.set("typeface", font_name)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    fill: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(emu_px(x), emu_px(y), emu_px(w), emu_px(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, FONT)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = LINE,
    radius: bool = True,
    line_width: float = 1.2,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, emu_px(x), emu_px(y), emu_px(w), emu_px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(line_width)
    return shape


def add_circle(slide, x: float, y: float, d: float, fill: str, label: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu_px(x), emu_px(y), emu_px(d), emu_px(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(fill)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_run_font(run, FONT)
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = color(WHITE)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, fill: str = "6B7785") -> None:
    head = 14
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu_px(x1), emu_px(y1), emu_px(x2 - head), emu_px(y2))
    conn.line.color.rgb = color(fill)
    conn.line.width = Pt(2.2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, emu_px(x2 - head), emu_px(y2 - head / 2), emu_px(head), emu_px(head))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = color(fill)
    tri.line.color.rgb = color(fill)


def add_step(slide, x: float, y: float, num: str, title: str, body: str, accent: str) -> None:
    add_rect(slide, x, y, 392, 166)
    add_circle(slide, x + 28, y + 32, 54, accent, num)
    add_text(slide, x + 106, y + 26, 240, 40, title, 22, accent, True)
    add_text(slide, x + 106, y + 82, 250, 58, body, 14, MUTED, True)


def add_metric(slide, x: float, y: float, value: str, label: str, accent: str) -> None:
    add_text(slide, x, y, 190, 54, value, 30, accent, True)
    add_text(slide, x + 180, y + 10, 270, 36, label, 15, MUTED, True)


def build_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, BG, BG, False, 0)

    add_text(slide, 86, 56, 320, 38, "PROJECTS", 15, MUTED, True)
    add_text(slide, 86, 112, 1250, 78, "증권사 리포트 텍스트 기반 산업 모니터링", 34, MUTED, True)
    add_text(slide, 1380, 136, 360, 36, "2024.05~2024.07", 17, MUTED, True)

    steps = [
        ("1", "수집", "증권사 기업 리포트 12.8만건\n기업·업종·지역 메타데이터", BLUE),
        ("2", "정제", "PDF 텍스트 추출\n수치·중복·공시문 제거", TEAL),
        ("3", "분석", "BERT 감성분석\nTrigram 이슈 키워드", PURPLE),
        ("4", "검증", "GDP·BSI·주가지수 비교\n선행성·예측력 점검", RED),
    ]
    x = 86
    for idx, args in enumerate(steps):
        add_step(slide, x, 238, *args)
        if idx < 3:
            add_arrow(slide, x + 404, 321, x + 450, 321)
        x += 456

    add_rect(slide, 92, 475, 760, 382)
    add_text(slide, 130, 512, 260, 44, "텍스트 처리 로직", 23, INK, True)
    add_text(slide, 130, 576, 145, 32, "문장 단위", 15, BLUE, True)
    add_text(slide, 300, 576, 420, 34, "리포트 본문 → 유효 문장 샘플", 15, INK, True)
    add_text(slide, 130, 638, 145, 32, "감성분석", 15, PURPLE, True)
    add_rect(slide, 300, 626, 104, 44, LIGHT, BLUE, False)
    add_text(slide, 300, 637, 104, 26, "긍정", 13, BLUE, True, PP_ALIGN.CENTER)
    add_rect(slide, 416, 626, 104, 44, LIGHT, RED, False)
    add_text(slide, 416, 637, 104, 26, "부정", 13, RED, True, PP_ALIGN.CENTER)
    add_rect(slide, 532, 626, 104, 44, LIGHT, LINE, False)
    add_text(slide, 532, 637, 104, 26, "중립", 13, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, 130, 704, 145, 32, "지수화", 15, TEAL, True)
    add_text(slide, 300, 696, 430, 40, "TBCI = (긍정 - 부정) / 감성 문장", 14, INK, True)
    add_text(slide, 130, 770, 145, 32, "요인 추출", 15, GOLD, True)
    add_text(slide, 300, 762, 430, 48, "요인·평가 분해 후 Trigram 빈도 집계", 14, INK, True)

    add_rect(slide, 925, 475, 430, 382)
    add_text(slide, 965, 512, 260, 44, "산출 지표", 23, INK, True)
    outputs = [
        ("TBCI", "업종별 텍스트 업황"),
        ("BC-factors", "분기별 경영환경 변화요인"),
        ("TIEI / TEEI", "이벤트 영향도와 평가"),
        ("BC-similarity", "공통요인 기반 산업 유사도"),
    ]
    y = 580
    for code, label in outputs:
        add_rect(slide, 965, y, 150, 42, LIGHT, LIGHT, False)
        add_text(slide, 980, y + 9, 120, 24, code, 12, BLUE, True, PP_ALIGN.CENTER)
        add_text(slide, 1140, y + 7, 180, 30, label, 13, MUTED, True)
        y += 62

    add_rect(slide, 1415, 475, 392, 382)
    add_text(slide, 1455, 512, 300, 44, "구현 포인트", 23, INK, True)
    add_metric(slide, 1455, 585, "1.45M", "유효 문장 규모", GREEN)
    add_metric(slide, 1455, 665, "40+", "산업 분류 매핑", BLUE)
    add_metric(slide, 1455, 745, "0.91", "선행지수 상관 예시", GOLD)

    add_rect(slide, 118, 910, 500, 82)
    add_text(slide, 150, 928, 170, 30, "자동화", 16, BLUE, True)
    add_text(slide, 315, 928, 250, 42, "수집·정제·분석·시각화 파이프라인", 14, MUTED, True)
    add_rect(slide, 710, 910, 500, 82)
    add_text(slide, 742, 928, 170, 30, "설명가능성", 16, TEAL, True)
    add_text(slide, 907, 928, 250, 42, "키워드와 원문 문장으로 결과 근거 확인", 14, MUTED, True)
    add_rect(slide, 1302, 910, 500, 82)
    add_text(slide, 1334, 928, 170, 30, "검증", 16, RED, True)
    add_text(slide, 1499, 928, 250, 42, "거시지표와 시차상관·Granger 검정", 14, MUTED, True)


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_slide(prs)
    out = OUT / "ai_industry_monitoring_editable.pptx"
    prs.save(out)
    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
